"""Generate arxiv_search.pptx — dark technical deck focused on the 4 NoSQL DBs."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# --- palette (matches the project's dark theme) ---
BG = RGBColor(0x0F, 0x11, 0x17)
SURFACE = RGBColor(0x1A, 0x1D, 0x27)
BORDER = RGBColor(0x2A, 0x2D, 0x3E)
ACCENT = RGBColor(0x63, 0x66, 0xF1)
ACCENT_SOFT = RGBColor(0x4F, 0x46, 0xE5)
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def hex2rgb(c: str) -> RGBColor:
    c = c.lstrip("#")
    return RGBColor(int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))


def add_blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=16, color=TEXT, bullet_color=ACCENT,
                line_spacing=1.2, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.name = font
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        r1.font.bold = True
        # text
        r2 = p.add_run()
        r2.text = line
        r2.font.name = font
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_card(slide, x, y, w, h, *, fill=SURFACE, border=BORDER, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    if radius:
        card.adjustments[0] = 0.06
    # remove default text
    card.text_frame.text = ""
    return card


def add_chip(slide, x, y, label, *, fill=ACCENT, text_color=RGBColor(0xFF, 0xFF, 0xFF),
             size=11, padding_x=0.18, padding_y=0.06):
    # estimate width by char count
    w = Inches(0.18 * len(label) + 2 * padding_x)
    h = Inches(0.32)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.fill.background()
    card.shadow.inherit = False
    card.adjustments[0] = 0.5
    tf = card.text_frame
    tf.margin_left = Inches(padding_x)
    tf.margin_right = Inches(padding_x)
    tf.margin_top = Inches(padding_y)
    tf.margin_bottom = Inches(padding_y)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.name = "Inter"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return card, w


def add_header(slide, kicker, title, *, kicker_color=ACCENT):
    add_text(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.35),
             kicker.upper(), size=11, bold=True, color=kicker_color)
    add_text(slide, Inches(0.7), Inches(0.78), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=TEXT)
    # divider
    line = slide.shapes.add_connector(1, Inches(0.7), Inches(1.55),
                                      Inches(13.333 - 0.7), Inches(1.55))
    line.line.color.rgb = BORDER
    line.line.width = Pt(0.75)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
             "ArXiv Semantic Search · NoSQL stack", size=9, color=TEXT_DIM)
    add_text(slide, Inches(11.633), Inches(7.05), Inches(1.0), Inches(0.3),
             f"{page_num} / {total}", size=9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def add_code(slide, x, y, w, h, code, *, size=12):
    card = add_card(slide, x, y, w, h, fill=RGBColor(0x10, 0x12, 0x18), border=BORDER)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15),
                                  w - Inches(0.4), h - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    lines = code.strip("\n").split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = line
        r.font.name = "JetBrains Mono"
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    return card


# ============================================================
# slide builders
# ============================================================

def slide_title(prs, total):
    s = add_blank_slide(prs)

    # accent block (left bar)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    # decorative dotted vector (faux network)
    add_text(s, Inches(0.9), Inches(2.2), Inches(12), Inches(0.5),
             "УЧЕБНЫЙ ПРОЕКТ · АНАЛИЗ ДАННЫХ", size=12, bold=True, color=ACCENT)
    add_text(s, Inches(0.9), Inches(2.7), Inches(12), Inches(1.4),
             "ArXiv Semantic Search", size=54, bold=True, color=TEXT)
    add_text(s, Inches(0.9), Inches(4.0), Inches(12), Inches(0.7),
             "Polyglot persistence: четыре NoSQL базы данных в одном проекте",
             size=22, color=TEXT_DIM)

    # tech chips
    chips = ["MongoDB", "Qdrant", "Redis", "ClickHouse", "FastAPI", "React"]
    cx = Inches(0.9)
    for label in chips:
        _, w = add_chip(s, cx, Inches(5.2), label,
                        fill=SURFACE, text_color=TEXT, size=11)
        cx += w + Inches(0.12)

    add_text(s, Inches(0.9), Inches(6.6), Inches(12), Inches(0.4),
             "РУДН  ·  Анализ данных", size=11, color=TEXT_DIM)


def slide_what(prs, page, total):
    s = add_blank_slide(prs)
    add_header(s, "Что за проект", "Семантический поиск по статьям arXiv")

    add_bullets(s, Inches(0.7), Inches(2.0), Inches(7.8), Inches(4.5), [
        "База ~8 000 статей категорий cs.AI / cs.LG / cs.CV",
        "Запрос на естественном языке, а не keyword-match",
        "Пример: «attention mechanism for image segmentation»",
        "Ответ — топ статей по смысловой близости + score",
        "Аналитика: какие запросы популярны, как растёт корпус по годам",
    ], size=17)

    # right panel — example query
    add_card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.5))
    add_text(s, Inches(9.1), Inches(2.15), Inches(3.5), Inches(0.4),
             "ПРИМЕР ОТВЕТА", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(9.1), Inches(2.55), Inches(3.5), Inches(0.5),
             "vision transformer", size=15, bold=True, color=TEXT)
    add_text(s, Inches(9.1), Inches(3.05), Inches(3.5), Inches(0.4),
             "topK = 10 · 42 ms · cached: false", size=10, color=TEXT_DIM)

    items = [
        ("Swin Transformer: Hierarchical…", "0.91"),
        ("ViT for Image Recognition at Scale", "0.88"),
        ("DeiT: Data-efficient Image…", "0.84"),
        ("Tokens-to-Token ViT", "0.81"),
    ]
    y = 3.55
    for title, score in items:
        add_text(s, Inches(9.1), Inches(y), Inches(2.6), Inches(0.35),
                 title, size=10, color=TEXT)
        add_text(s, Inches(11.7), Inches(y), Inches(0.9), Inches(0.35),
                 score, size=10, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
        y += 0.55

    add_footer(s, page, total)


def slide_why_four(prs, page, total):
    s = add_blank_slide(prs)
    add_header(s, "Почему 4 разных БД", "Polyglot persistence — каждая БД делает то, что умеет лучше всех")

    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
             "Один тип запросов одна БД делает идеально, остальные — плохо. Поэтому в одном сервисе сосуществуют 4 разных движка:",
             size=14, color=TEXT_DIM)

    # 4 cards row
    cards = [
        ("MongoDB", "Метаданные", "Document store", "хранит JSON-объекты статей со\nсложной вложенной структурой"),
        ("Qdrant", "Векторный поиск", "Vector DB", "ANN-поиск по эмбеддингам,\nкосинусная близость"),
        ("Redis", "Кэш", "In-memory KV", "повторный запрос — ответ\nза доли миллисекунды из RAM"),
        ("ClickHouse", "Аналитика", "Columnar OLAP", "быстрые GROUP BY и агрегации\nна логах поисков"),
    ]
    card_w = Inches(2.95)
    gap = Inches(0.15)
    x = Inches(0.7)
    for name, role, kind, desc in cards:
        add_card(s, x, Inches(2.9), card_w, Inches(3.6))
        add_text(s, x + Inches(0.25), Inches(3.05), card_w - Inches(0.5), Inches(0.4),
                 role.upper(), size=10, bold=True, color=ACCENT)
        add_text(s, x + Inches(0.25), Inches(3.45), card_w - Inches(0.5), Inches(0.5),
                 name, size=20, bold=True, color=TEXT)
        add_text(s, x + Inches(0.25), Inches(4.05), card_w - Inches(0.5), Inches(0.4),
                 kind, size=11, color=TEXT_DIM)
        # divider
        ln = s.shapes.add_connector(1, x + Inches(0.25), Inches(4.55),
                                    x + card_w - Inches(0.25), Inches(4.55))
        ln.line.color.rgb = BORDER
        ln.line.width = Pt(0.5)
        add_text(s, x + Inches(0.25), Inches(4.7), card_w - Inches(0.5), Inches(1.5),
                 desc, size=11, color=TEXT)
        x += card_w + gap

    add_footer(s, page, total)


def slide_architecture(prs, page, total):
    s = add_blank_slide(prs)
    add_header(s, "Архитектура", "Как сервисы связаны между собой")

    # Frontend node
    def node(x, y, w, h, label, sub, fill=SURFACE):
        add_card(s, x, y, w, h, fill=fill)
        add_text(s, x, y + Inches(0.18), w, Inches(0.4), label,
                 size=14, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.55), w, Inches(0.35), sub,
                 size=10, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    # Frontend
    node(Inches(0.9), Inches(3.2), Inches(2.2), Inches(1.0),
         "Frontend", "React + Vite · :5173")
    # Backend
    node(Inches(4.9), Inches(3.2), Inches(2.6), Inches(1.0),
         "Backend", "FastAPI · :8000", fill=RGBColor(0x21, 0x24, 0x35))
    # 4 DBs (right column)
    db_x = Inches(9.5)
    db_y = Inches(1.9)
    db_w = Inches(3.2)
    db_h = Inches(0.85)
    dbs = [
        ("MongoDB", "метаданные :27017", RGBColor(0x10, 0x4E, 0x35)),
        ("Qdrant", "векторы :6333", RGBColor(0x3B, 0x2A, 0x6B)),
        ("Redis", "кэш :6379", RGBColor(0x6B, 0x1F, 0x1F)),
        ("ClickHouse", "OLAP :8123", RGBColor(0x6B, 0x4F, 0x1F)),
    ]
    for i, (n, sub, col) in enumerate(dbs):
        y = db_y + Emu(int(db_h) * i + Inches(0.15).emu * i)
        add_card(s, db_x, y, db_w, db_h, fill=col, border=BORDER)
        add_text(s, db_x + Inches(0.2), y + Inches(0.12), db_w, Inches(0.35),
                 n, size=13, bold=True, color=TEXT)
        add_text(s, db_x + Inches(0.2), y + Inches(0.44), db_w, Inches(0.35),
                 sub, size=10, color=TEXT_DIM)

    # arrows
    def arrow(x1, y1, x2, y2):
        ln = s.shapes.add_connector(1, x1, y1, x2, y2)  # straight
        ln.line.color.rgb = ACCENT
        ln.line.width = Pt(1.5)

    arrow(Inches(3.1), Inches(3.7), Inches(4.9), Inches(3.7))
    # backend → 4 DBs
    bx = Inches(7.5)
    by = Inches(3.7)
    for i in range(4):
        ty = db_y + Emu(int(db_h) * i + Inches(0.15).emu * i) + Inches(0.42)
        arrow(bx, by, db_x, ty)

    add_text(s, Inches(0.7), Inches(5.6), Inches(12), Inches(0.4),
             "На каждый запрос /api/search бэк ходит во все 4 БД (по очереди / параллельно).",
             size=12, color=TEXT_DIM)
    add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
             "Развёртывание: один docker compose up — 6 контейнеров с healthcheck-зависимостями.",
             size=12, color=TEXT_DIM)

    add_footer(s, page, total)


def slide_db_theory(prs, page, total, *, kicker, name, subtitle, theory_bullets,
                    inner_label, inner_text):
    s = add_blank_slide(prs)
    add_header(s, kicker, f"{name} — что это и как устроено")

    add_text(s, Inches(0.7), Inches(2.0), Inches(8), Inches(0.5),
             subtitle, size=14, color=TEXT_DIM)
    add_bullets(s, Inches(0.7), Inches(2.7), Inches(7.5), Inches(4.0),
                theory_bullets, size=15)

    add_card(s, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.6))
    add_text(s, Inches(9.1), Inches(2.15), Inches(3.5), Inches(0.4),
             "ВНУТРЕННЕЕ УСТРОЙСТВО", size=10, bold=True, color=ACCENT)
    add_text(s, Inches(9.1), Inches(2.55), Inches(3.5), Inches(0.4),
             inner_label, size=14, bold=True, color=TEXT)
    add_text(s, Inches(9.1), Inches(3.05), Inches(3.5), Inches(3.5),
             inner_text, size=11, color=TEXT)

    add_footer(s, page, total)


def slide_db_usage(prs, page, total, *, kicker, name, where, how_bullets, code, code_lang_label):
    s = add_blank_slide(prs)
    add_header(s, kicker, f"{name} — как использовано в проекте")

    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
             where, size=14, color=TEXT_DIM)

    add_bullets(s, Inches(0.7), Inches(2.7), Inches(6.3), Inches(4.5),
                how_bullets, size=14)

    add_text(s, Inches(7.4), Inches(2.7), Inches(5.5), Inches(0.4),
             code_lang_label, size=10, bold=True, color=ACCENT)
    add_code(s, Inches(7.4), Inches(3.05), Inches(5.3), Inches(3.7), code, size=11)

    add_footer(s, page, total)


def slide_request_flow(prs, page, total):
    s = add_blank_slide(prs)
    add_header(s, "End-to-end", "Полный путь одного запроса /api/search")

    steps = [
        ("1", "Redis", "GET search:{query}", "промах → идём дальше", RGBColor(0x6B, 0x1F, 0x1F)),
        ("2", "Encoder", "embed(query) → 384-dim", "all-MiniLM-L6-v2", RGBColor(0x21, 0x24, 0x35)),
        ("3", "Qdrant", "ANN top-K по cosine", "→ список arxiv_id + score", RGBColor(0x3B, 0x2A, 0x6B)),
        ("4", "MongoDB", "find({arxiv_id: {$in:[…]}})", "достаём метаданные", RGBColor(0x10, 0x4E, 0x35)),
        ("5", "Redis", "SETEX 3600 JSON(payload)", "кладём в кэш на час", RGBColor(0x6B, 0x1F, 0x1F)),
        ("6", "ClickHouse", "INSERT INTO search_logs", "лог запроса для аналитики", RGBColor(0x6B, 0x4F, 0x1F)),
        ("7", "Response", "JSON → клиент", "результат + query_time_ms", RGBColor(0x21, 0x24, 0x35)),
    ]

    y = Inches(2.0)
    row_h = Inches(0.62)
    gap = Inches(0.08)
    for num, where, op, note, color in steps:
        # number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.1),
                                   Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        badge.shadow.inherit = False
        tf = badge.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = "Inter"; r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # row card
        add_card(s, Inches(1.3), y, Inches(11.4), row_h, fill=SURFACE)
        # where chip
        chip, w = add_chip(s, Inches(1.5), y + Inches(0.17), where,
                           fill=color, text_color=TEXT, size=10)
        # op (mono)
        op_tb = s.shapes.add_textbox(Inches(1.5) + w + Inches(0.2),
                                     y + Inches(0.15), Inches(5.5), Inches(0.4))
        op_tf = op_tb.text_frame
        op_tf.margin_left = Emu(0); op_tf.margin_right = Emu(0)
        op_tf.margin_top = Emu(0); op_tf.margin_bottom = Emu(0)
        p = op_tf.paragraphs[0]
        r = p.add_run()
        r.text = op
        r.font.name = "JetBrains Mono"
        r.font.size = Pt(11)
        r.font.color.rgb = TEXT
        # note
        add_text(s, Inches(8.0), y + Inches(0.18), Inches(4.5), Inches(0.4),
                 note, size=11, color=TEXT_DIM)

        y = y + row_h + gap

    add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
             "Кэш-хит:  шаги 1 → 6 → 7  (≈ 5 ms).  Полный путь:  ≈ 40-80 ms на тёплой модели.",
             size=12, color=TEXT_DIM, align=PP_ALIGN.LEFT)

    add_footer(s, page, total)


def slide_compare(prs, page, total):
    s = add_blank_slide(prs)
    add_header(s, "Итог", "Сравнение всех четырёх БД в проекте")

    headers = ["БД", "Тип", "Модель данных", "Роль в проекте", "Ключевая операция"]
    rows = [
        ("MongoDB", "Document", "BSON-документы", "Метаданные статей", "find({arxiv_id: {$in: …}})"),
        ("Qdrant", "Vector", "Точки + 384-d векторы", "Семантический поиск", "ANN search (cosine, HNSW)"),
        ("Redis", "Key-Value", "Строки в RAM", "Кэш ответов", "GET / SETEX 3600"),
        ("ClickHouse", "Columnar OLAP", "Колонки с сжатием", "Логи и статистика", "INSERT + GROUP BY"),
    ]

    col_widths = [Inches(1.7), Inches(1.6), Inches(2.6), Inches(2.8), Inches(3.3)]
    x0 = Inches(0.7)
    y0 = Inches(2.1)
    row_h = Inches(0.55)

    # header row
    x = x0
    for i, h in enumerate(headers):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, col_widths[i], row_h)
        cell.fill.solid(); cell.fill.fore_color.rgb = SURFACE
        cell.line.color.rgb = BORDER; cell.line.width = Pt(0.5)
        cell.shadow.inherit = False
        cell.text_frame.text = ""
        add_text(s, x + Inches(0.15), y0 + Inches(0.13), col_widths[i] - Inches(0.3), Inches(0.35),
                 h.upper(), size=10, bold=True, color=ACCENT)
        x += col_widths[i]

    # body rows
    for r_idx, row in enumerate(rows):
        y = y0 + row_h + Emu(int(row_h) * r_idx)
        x = x0
        for i, val in enumerate(row):
            cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_widths[i], row_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG if r_idx % 2 == 0 else RGBColor(0x14, 0x16, 0x1E)
            cell.line.color.rgb = BORDER; cell.line.width = Pt(0.5)
            cell.shadow.inherit = False
            cell.text_frame.text = ""
            font = "JetBrains Mono" if i == 4 else "Inter"
            size = 11 if i == 4 else 12
            bold = i == 0
            add_text(s, x + Inches(0.15), y + Inches(0.16),
                     col_widths[i] - Inches(0.3), Inches(0.35),
                     val, size=size, bold=bold, color=TEXT, font=font)
            x += col_widths[i]

    add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
             "Каждая БД оптимизирована под свой класс задач. Заменить любую на другую — деградация по скорости или функциональности.",
             size=12, color=TEXT_DIM)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.4),
             "Это и есть polyglot persistence — стандартный подход в современных production-системах.",
             size=12, color=TEXT_DIM)

    add_footer(s, page, total)


def slide_outro(prs, page, total):
    s = add_blank_slide(prs)
    add_header(s, "Источники и ссылки", "Документация использованных БД")

    items = [
        ("MongoDB", "https://www.mongodb.com/docs/manual/"),
        ("Qdrant", "https://qdrant.tech/documentation/"),
        ("Redis", "https://redis.io/docs/latest/"),
        ("ClickHouse", "https://clickhouse.com/docs"),
        ("sentence-transformers", "https://www.sbert.net/"),
        ("HNSW (Hierarchical NSW)", "https://arxiv.org/abs/1603.09320"),
        ("Polyglot persistence (M. Fowler)", "https://martinfowler.com/bliki/PolyglotPersistence.html"),
    ]
    y = 2.05
    for label, url in items:
        add_text(s, Inches(0.7), Inches(y), Inches(4.0), Inches(0.4),
                 label, size=14, bold=True, color=TEXT)
        add_text(s, Inches(4.8), Inches(y), Inches(8.0), Inches(0.4),
                 url, size=12, color=ACCENT, font="JetBrains Mono")
        y += 0.5

    add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.4),
             "Спасибо!", size=22, bold=True, color=TEXT)

    add_footer(s, page, total)


# ============================================================
# build
# ============================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Total = 14 slides
    total = 14
    page = 0

    slide_title(prs, total); page += 1
    slide_what(prs, page + 1 - 1 + 1, total); page += 1   # 2

    # 3
    page += 1
    slide_why_four(prs, page, total)
    # 4
    page += 1
    slide_architecture(prs, page, total)

    # MongoDB pair
    page += 1
    slide_db_theory(
        prs, page, total,
        kicker="БД 1/4 · MongoDB",
        name="MongoDB",
        subtitle="Документная NoSQL БД. Хранит JSON-подобные объекты (BSON) вместо строк таблицы.",
        theory_bullets=[
            "Без жёсткой схемы — каждый документ может иметь свой набор полей",
            "Удобно для сущностей с массивами и вложенными объектами (authors[], categories[])",
            "Запросы на JSON-подобном языке: find, aggregate (pipeline)",
            "Индексы B-tree на полях, как в SQL — быстрый lookup по arxiv_id",
            "Горизонтально масштабируется через шардинг",
        ],
        inner_label="BSON + B-tree индексы",
        inner_text=(
            "Документы хранятся в формате BSON (binary JSON).\n\n"
            "Для каждого индекса MongoDB строит B-tree по значениям поля.\n\n"
            "Поиск по индексированному полю — O(log n), как в реляционных БД.\n\n"
            "Aggregation pipeline — последовательность стадий ($match, $group, $sort) "
            "над коллекцией, аналог SQL GROUP BY с гибкой композицией."
        ),
    )
    page += 1
    slide_db_usage(
        prs, page, total,
        kicker="БД 1/4 · MongoDB",
        name="MongoDB",
        where="Коллекция articles в БД arxiv. Полные метаданные статьи + индекс на arxiv_id.",
        how_bullets=[
            "Загрузка — bulk_write с upsert по arxiv_id (идемпотентно)",
            "Поиск — find({arxiv_id: {$in: [...]}}), результат top-K из Qdrant",
            "Аналитика по годам и категориям — aggregation pipeline",
            "Один индекс delivers быстрый lookup, второй — диапазон по годам",
        ],
        code_lang_label="ДОКУМЕНТ + ЗАПРОС",
        code=(
            "// один документ в коллекции articles\n"
            "{\n"
            '  "arxiv_id":   "2301.00001",\n'
            '  "title":      "Attention Is All You Need",\n'
            '  "authors":    ["A. Vaswani", "N. Shazeer"],\n'
            '  "year":       2017,\n'
            '  "categories": ["cs.LG", "cs.CL"],\n'
            '  "abstract":   "..."\n'
            "}\n\n"
            "// агрегация по годам\n"
            "db.articles.aggregate([\n"
            '  { $group: { _id: "$year", count: { $sum: 1 } } },\n'
            "  { $sort:  { _id: 1 } }\n"
            "])"
        ),
    )

    # Qdrant pair
    page += 1
    slide_db_theory(
        prs, page, total,
        kicker="БД 2/4 · Qdrant",
        name="Qdrant",
        subtitle="Векторная БД. Хранит точки в N-мерном пространстве и ищет ближайшие соседи.",
        theory_bullets=[
            "Каждая запись = вектор + payload (метаданные)",
            "Эмбеддинг текста: модель преобразует абстракт в 384-мерный вектор",
            "Похожие по смыслу тексты → близкие точки в пространстве",
            "Метрика — cosine similarity (угол между векторами)",
            "Точный k-NN на миллионах точек слишком медленный → ANN-индекс",
        ],
        inner_label="HNSW (Hierarchical NSW)",
        inner_text=(
            "Hierarchical Navigable Small World — алгоритм approximate nearest neighbor.\n\n"
            "Многоуровневый граф: верхние слои разрежены и ускоряют переход к нужной "
            "области, нижние — плотные для точного поиска.\n\n"
            "Поиск ~ O(log n) с очень высокой recall (>0.95 при правильных параметрах).\n\n"
            "Trade-off: память на индекс vs скорость поиска."
        ),
    )
    page += 1
    slide_db_usage(
        prs, page, total,
        kicker="БД 2/4 · Qdrant",
        name="Qdrant",
        where="Коллекция arxiv: точки с 384-мерными векторами абстрактов и payload {arxiv_id, title, year, categories}.",
        how_bullets=[
            "При загрузке — encode(abstract) → upsert PointStruct в коллекцию",
            "При поиске — encode(query) → search(vector, limit=10)",
            "Ответ Qdrant: список (id, score, payload) — сразу есть arxiv_id",
            "Метрика Cosine, нормализуем эмбеддинги (normalize_embeddings=True)",
        ],
        code_lang_label="СОЗДАНИЕ И ПОИСК",
        code=(
            "# один раз на старте\n"
            "client.create_collection(\n"
            '    collection_name="arxiv",\n'
            "    vectors_config=VectorParams(\n"
            "        size=384, distance=Distance.COSINE,\n"
            "    ),\n"
            ")\n\n"
            "# поисковый запрос\n"
            "vec = model.encode(query, normalize_embeddings=True)\n"
            "hits = client.search(\n"
            '    collection_name="arxiv",\n'
            "    query_vector=vec.tolist(),\n"
            "    limit=10, with_payload=True,\n"
            ")"
        ),
    )

    # Redis pair
    page += 1
    slide_db_theory(
        prs, page, total,
        kicker="БД 3/4 · Redis",
        name="Redis",
        subtitle="In-memory key-value хранилище. Всё в RAM, поэтому миллионы операций в секунду.",
        theory_bullets=[
            "Структуры: строки, хеши, списки, множества, sorted sets, streams",
            "Однопоточный event loop — нет конкуренции на структурах данных",
            "TTL: ключ автоматически удаляется через N секунд",
            "Опциональная персистентность (RDB / AOF), но это не основной use-case",
            "Типичные роли: кэш, очередь, rate limiter, sessions",
        ],
        inner_label="Hash table + TTL bucket",
        inner_text=(
            "Все ключи лежат в большой hash table в RAM.\n\n"
            "Чтение/запись — O(1) для строк, амортизированно.\n\n"
            "TTL хранится отдельно. Ленивое удаление при доступе + фоновый sweep "
            "случайных ключей с истёкшим временем.\n\n"
            "Для нашего сценария важно: ответ из RAM в десятки раз быстрее, "
            "чем повторный полный путь embed → ANN → mongo."
        ),
    )
    page += 1
    slide_db_usage(
        prs, page, total,
        kicker="БД 3/4 · Redis",
        name="Redis",
        where="Кэш ответов /api/search. Ключ — нормализованный текст запроса, значение — JSON ответа.",
        how_bullets=[
            "Ключ: search:{query.lower().strip()} — нормализуем перед хешем",
            "Значение: сериализованный JSON-payload результата",
            "TTL = 1 час (SETEX) — данные слабо меняются, кэш можно держать",
            "На каждом /search сначала GET, при промахе — основная логика",
            "Хит сокращает latency с ~50 мс до ~3 мс",
        ],
        code_lang_label="ЛОГИКА КЭША",
        code=(
            "def get_cached(query: str):\n"
            '    key = f"search:{normalize(query)}"\n'
            "    raw = r.get(key)\n"
            "    return json.loads(raw) if raw else None\n\n"
            "def set_cached(query, payload):\n"
            '    key = f"search:{normalize(query)}"\n'
            "    r.setex(key, 3600, json.dumps(payload))\n\n"
            "# в обработчике\n"
            "if cached := get_cached(q):\n"
            "    return {**cached, \"cached\": True}"
        ),
    )

    # ClickHouse pair
    page += 1
    slide_db_theory(
        prs, page, total,
        kicker="БД 4/4 · ClickHouse",
        name="ClickHouse",
        subtitle="Колоночная OLAP-СУБД. Заточена под аналитические запросы по миллиардам строк.",
        theory_bullets=[
            "Хранит данные по колонкам, а не по строкам — читаем только нужные поля",
            "Высокое сжатие (LZ4/ZSTD) — однотипные значения в колонке",
            "Векторизованное выполнение запросов (SIMD)",
            "Стандартный SQL + расширения (массивы, semi-structured)",
            "Оптимизирована под INSERT-batch и SELECT-aggregate, не под точечные UPDATE",
        ],
        inner_label="MergeTree engine",
        inner_text=(
            "Семейство движков MergeTree — основа ClickHouse.\n\n"
            "Данные пишутся в неизменяемые «парты» (parts), отсортированные по "
            "ORDER BY ключу.\n\n"
            "Фоновое слияние (merge) объединяет мелкие парты в крупные → "
            "поддерживает компактность и скорость чтения.\n\n"
            "На SELECT движок читает только нужные колонки и пропускает блоки "
            "по индексу разрежённой сортировки (sparse primary index)."
        ),
    )
    page += 1
    slide_db_usage(
        prs, page, total,
        kicker="БД 4/4 · ClickHouse",
        name="ClickHouse",
        where="Таблица search_logs хранит каждый поисковый запрос. Источник данных для дашборда аналитики.",
        how_bullets=[
            "Каждый /search → INSERT в search_logs (query, results_count, timestamp)",
            "Топ-10 запросов: GROUP BY query ORDER BY count() DESC",
            "Активность по дням / общее число поисков — toDate(timestamp), count()",
            "На фронте таб Analytics дёргает /api/stats, бэк делает SQL в ClickHouse",
        ],
        code_lang_label="ТАБЛИЦА И ЗАПРОСЫ",
        code=(
            "CREATE TABLE search_logs (\n"
            "    query         String,\n"
            "    results_count UInt32,\n"
            "    timestamp     DateTime DEFAULT now()\n"
            ") ENGINE = MergeTree()\n"
            "ORDER BY timestamp;\n\n"
            "-- top queries\n"
            "SELECT query, count() AS cnt\n"
            "FROM search_logs\n"
            "GROUP BY query\n"
            "ORDER BY cnt DESC\n"
            "LIMIT 10;"
        ),
    )

    # request flow + compare + outro
    page += 1
    slide_request_flow(prs, page, total)
    page += 1
    slide_compare(prs, page, total)
    page += 1
    slide_outro(prs, page, total)

    out = "arxiv_search.pptx"
    prs.save(out)
    print(f"saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
